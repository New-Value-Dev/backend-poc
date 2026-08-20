"""PPTX(파워포인트) 파서"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.table import Table as PptxTable

from app.services.parsing.base import BODY_LEVEL, UnifiedDocument, UnifiedSection

_TITLE_PLACEHOLDER_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


def _is_title_shape(shape) -> bool:
    # placeholder가 아닌 도형(텍스트박스/이미지 등)에서 placeholder_format 에 접근하면
    # None이 아니라 ValueError를 던지므로, is_placeholder로 먼저 걸러야 한다.
    if not shape.is_placeholder:
        return False
    try:
        return shape.placeholder_format.type in _TITLE_PLACEHOLDER_TYPES
    except ValueError:
        # 일부 레이아웃 전용 placeholder 타입은 python-pptx 가 매핑을 못 해 예외를 던진다.
        return False


def table_text(table: PptxTable) -> str:
    rows = []
    for row in table.rows:
        rows.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(rows)


def _iter_shapes(shapes):
    """그룹 도형은 그 자체로 has_table/has_text_frame이 없으므로, 그룹이면 안의 도형들을
    재귀적으로 풀어서 내려준다 (그룹 중첩도 지원)."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def iter_slide_blocks(prs: Presentation):
    """(kind, obj, text, slide_index) 를 슬라이드 -> 도형 순서대로 yield.

    kind: "heading"(제목 placeholder) | "paragraph"(본문 문단) | "table". obj 는
    correction_service 가 다시 찾아 고쳐 쓸 수 있도록 문단/표 객체 그대로 준다
    (문단은 run.text getter/setter, 표는 셀을 순회하며 원문이 있는 셀만 찾아 고친다).
    """
    for slide_index, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            if shape.has_table:
                text = table_text(shape.table)
                if text.strip():
                    yield "table", shape.table, text, slide_index
                continue
            if not shape.has_text_frame:
                continue
            is_title = _is_title_shape(shape)
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                yield ("heading" if is_title else "paragraph"), paragraph, text, slide_index


class PptxParser:
    name = "pptx"
    extensions = ("pptx",)

    def parse(self, path: Path) -> UnifiedDocument:
        prs = Presentation(str(path))
        sections: list[UnifiedSection] = []
        order = 0

        for kind, _obj, text, slide_index in iter_slide_blocks(prs):
            meta = {"slide_index": slide_index}
            if kind == "heading":
                sections.append(
                    UnifiedSection(
                        content=text, title=text, order=order,
                        level=1, section_type="heading", meta=meta,
                    )
                )
            elif kind == "table":
                sections.append(
                    UnifiedSection(
                        content=text, order=order,
                        level=BODY_LEVEL, section_type="table", meta=meta,
                    )
                )
            else:
                sections.append(
                    UnifiedSection(
                        content=text, order=order,
                        level=BODY_LEVEL, section_type="paragraph", meta=meta,
                    )
                )
            order += 1

        return UnifiedDocument(
            sections=sections, parser_name=self.name,
            meta={"slide_count": len(prs.slides)},
        )
